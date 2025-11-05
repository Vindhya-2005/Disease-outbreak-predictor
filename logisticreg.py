from dotenv import load_dotenv
load_dotenv()
from flask import Flask, render_template, request, jsonify
import pandas as pd
import numpy as np
import folium
from sklearn.preprocessing import LabelEncoder, StandardScaler
from sklearn.linear_model import LogisticRegression
from sklearn.metrics import accuracy_score
import matplotlib.cm as cm
import matplotlib.colors as colors
import os
from pptx import Presentation
from langchain_core.documents import Document
from langchain_chroma import Chroma
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from mistralai.client import MistralClient

# --------------------------
# FLASK APP
# --------------------------
app = Flask(__name__)

# ==========================
# PART 1: ML MODEL SECTION
# ==========================
df = pd.read_excel("mlproject.xlsx")

# Data cleaning
for col in ['preci', 'LAI', 'Temp', 'Cases', 'Deaths']:
    if col in df.columns:
        df[col] = pd.to_numeric(df[col], errors='coerce')
        df[col] = df[col].fillna(df[col].mean())

# Lag features
for col in ['preci', 'LAI', 'Temp']:
    df[f'{col}_lag1'] = df.groupby('district')[col].shift(1)
    df[f'{col}_lag2'] = df.groupby('district')[col].shift(2)
    df[f'{col}_lag3'] = df.groupby('district')[col].shift(3)

df.dropna(inplace=True)

# Encoding
le_district = LabelEncoder()
df['district_encoded'] = le_district.fit_transform(df['district'])
le_disease = LabelEncoder()
df['Disease_encoded'] = le_disease.fit_transform(df['Disease'])

features = [
    'district_encoded','day','mon','year','Latitude','Longitude',
    'preci','LAI','Temp','preci_lag1','LAI_lag1','Temp_lag1',
    'preci_lag2','LAI_lag2','Temp_lag2','preci_lag3','LAI_lag3','Temp_lag3'
]
features = [f for f in features if f in df.columns]

X = df[features]
y = df['Disease_encoded']

scaler = StandardScaler()
X_scaled = scaler.fit_transform(X)

# --- Logistic Regression Model ---
lr_model = LogisticRegression(max_iter=1000, multi_class='multinomial', solver='lbfgs')
lr_model.fit(X_scaled, y)

from sklearn.metrics import classification_report, confusion_matrix

# Evaluate model performance
y_pred = lr_model.predict(X_scaled)

accuracy = accuracy_score(y, y_pred)
print(f"Logistic Regression Model trained successfully! (Accuracy: {accuracy:.4f})")

print("\n--- Classification Report ---")
print(classification_report(y, y_pred, target_names=le_disease.classes_))

print("\n--- Confusion Matrix ---")
print(confusion_matrix(y, y_pred))

district_coords = df.groupby('district')[['Latitude','Longitude']].first().to_dict('index')

# Predict function
def predict_outbreaks_for_disease(year, month, disease_name):
    if disease_name not in le_disease.classes_:
        return pd.DataFrame()

    all_districts = df['district'].unique()
    avg_climate = df.groupby(['district','mon'])[['preci','LAI','Temp',
                     'preci_lag1','LAI_lag1','Temp_lag1',
                     'preci_lag2','LAI_lag2','Temp_lag2',
                     'preci_lag3','LAI_lag3','Temp_lag3']].mean().to_dict('index')

    future_rows = []
    for district in all_districts:
        key = (district, month)
        climate_data = avg_climate.get(key, avg_climate.get((district, 1), {})).copy()
        climate_data['Temp'] = climate_data.get('Temp',0) + (year - 2025)*0.2
        climate_data['preci'] = climate_data.get('preci',0) + (year - 2025)*1.0
        climate_data['LAI'] = climate_data.get('LAI',0) + (year - 2025)*0.05

        coords = district_coords.get(district, {'Latitude':0,'Longitude':0})
        future_rows.append({
            'district': district,
            'day': 15,
            'mon': month,
            'year': year,
            'Latitude': coords['Latitude'],
            'Longitude': coords['Longitude'],
            **climate_data
        })

    future_df = pd.DataFrame(future_rows)
    future_df['district_encoded'] = le_district.transform(future_df['district'])
    for col in features:
        if col in future_df.columns:
            future_df[col] = future_df[col].fillna(df[col].mean())

    future_X_scaled = scaler.transform(future_df[features])
    probs = lr_model.predict_proba(future_X_scaled)
    disease_index = np.where(lr_model.classes_ == le_disease.transform([disease_name])[0])[0][0]
    future_df['predicted_prob'] = probs[:, disease_index]
    future_df['predicted_cases'] = future_df['predicted_prob'] * 1000
    future_df['predicted_disease'] = disease_name
    return future_df[future_df['predicted_prob']>0.1]

def create_map(pred_df, disease_name, year, month):
    if pred_df.empty:
        return None
    cmap = cm.get_cmap("Reds")
    norm = colors.Normalize(vmin=pred_df['predicted_prob'].min(), vmax=pred_df['predicted_prob'].max())
    m = folium.Map(location=[22.5, 80], zoom_start=5)
    for _, row in pred_df.iterrows():
        rgb = cmap(norm(row['predicted_prob']))
        hex_color = colors.to_hex(rgb)
        radius = 5 + (row['predicted_cases']/200)**0.5
        folium.CircleMarker(
            location=[row['Latitude'], row['Longitude']],
            radius=radius,
            color=hex_color,
            fill=True,
            fill_color=hex_color,
            fill_opacity=0.9,
            popup=f"<b>{row['district']}</b>, {row['year']}-{row['mon']}<br>"
                  f"Disease: {row['predicted_disease']}<br>"
                  f"Probability: {row['predicted_prob']:.2f}"
        ).add_to(m)
    file_name = f"static/map_{disease_name}_{year}_{month}.html".replace(" ","")
    m.save(file_name)
    return file_name

# ==========================
# PART 2: CHATBOT SECTION
# ==========================
folder_path = "Data"
chroma_path = "chroma"

def load_documents():
    from langchain.document_loaders import PyPDFLoader
    docs = []
    for filename in os.listdir(folder_path):
        path = os.path.join(folder_path, filename)
        if filename.endswith(".pdf"):
            loader = PyPDFLoader(path)
            docs += loader.load()
        elif filename.endswith((".ppt", ".pptx")):
            pres = Presentation(path)
            text = ""
            for slide in pres.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text += shape.text + "\n"
            docs.append(Document(page_content=text, metadata={"source": path}))
    return docs

def prepare_chroma():
    if os.path.exists(chroma_path):
        return
    docs = load_documents()
    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    chunks = splitter.split_documents(docs)
    embeddings = HuggingFaceEmbeddings(model_name="all-mpnet-base-v2")
    Chroma.from_documents(chunks, embeddings, persist_directory=chroma_path)
    print(f"Knowledge base saved with {len(chunks)} chunks.")

def query_chatbot(query):
    embeddings = HuggingFaceEmbeddings(model_name="all-mpnet-base-v2")
    db = Chroma(persist_directory=chroma_path, embedding_function=embeddings)
    retriever = db.as_retriever(search_kwargs={'k': 4})
    relevant_docs = retriever.invoke(query)
    context = "\n\n".join([d.page_content for d in relevant_docs])

    prompt = f"""
Answer the question using the provided document text.
If found, cite "(from PDF)" or "(from PowerPoint)".
If not, use your general knowledge.

Question: {query}
Documents: {context}
Answer:
"""
    try:
        api_key = os.environ.get("MISTRAL_API_KEY")
        if not api_key:
            raise ValueError("MISTRAL_API_KEY not set.")

        client = MistralClient(api_key=api_key)
        response = client.chat(
            model="mistral-medium",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7
        )

        reply_text = response.choices[0].message.content
        return reply_text

    except Exception as e:
        return f"Error: {e}"

prepare_chroma()

# ==========================
# FLASK ROUTES
# ==========================
@app.route("/", methods=["GET","POST"])
def index():
    map_file = None
    disease_name = None
    precautions = []
    all_diseases = sorted(df['Disease'].unique())

    precautions_dict = {
        "Dengue":["Use mosquito nets","Avoid stagnant water","Wear long-sleeved clothes","Keep surroundings clean"],
        "Malaria":["Sleep under nets","Take preventive meds","Drain standing water","Spray indoors"],
        "Cholera":["Drink boiled water","Wash hands","Eat freshly cooked food","Avoid raw seafood"],
        "Acute Diarrhoeal Disease":["Drink safe water","Eat fresh food","Wash hands often","Use sanitary toilets"],
        "Acute Encephalatis Syndrome":["Seek care early","Maintain hygiene","Avoid mosquito bites","Sleep under nets"],
        "Other Diseases":["Maintain hygiene","Seek medical help promptly","Avoid crowded places"]
    }

    if request.method=="POST":
        year = int(request.form["year"])
        month = int(request.form["month"])
        disease_name = request.form["disease"]
        preds = predict_outbreaks_for_disease(year, month, disease_name)
        map_file = create_map(preds, disease_name, year, month)
        precautions = precautions_dict.get(disease_name, precautions_dict["Other Diseases"])

    return render_template("index.html", map_file=map_file, disease_name=disease_name,
                           precautions=precautions, all_diseases=all_diseases)

import markdown

@app.route("/chat", methods=["POST"])
def chat():
    user_msg = request.json.get("message", "")
    reply = query_chatbot(user_msg)
    html_reply = markdown.markdown(reply)
    return jsonify({"reply": html_reply})


if __name__ == "__main__":
    app.run(debug=True)
